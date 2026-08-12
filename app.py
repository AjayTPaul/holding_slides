"""
Informa Connect Session Slide Builder

Generate professional holding slides for events.
Usage: streamlit run app.py
"""

import base64
import io
import os
import subprocess
import tempfile
from copy import deepcopy
from pathlib import Path

import streamlit as st
from streamlit.components.v1 import html as st_html
import streamlit.components.v1 as components
import pandas as pd
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn


# -----------------------------------------------------------------------------
# CONFIGURATION
# -----------------------------------------------------------------------------

SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5
SLIDE_ASPECT = SLIDE_WIDTH_IN / SLIDE_HEIGHT_IN  # ~1.778 (16:9)

COLORS = {
    "title_text_dark": "283857",  # Carbon
    "title_text_light": "FFFFFF",  # White
    "speaker_text_dark": "283857",
    "speaker_text_light": "FFFFFF",
    "default_bg": "F5F5F5",
    "cta_bg": "002244",  # Indigo
}

CONTENT_PANEL = {
    "left": 1.5,
    "top": 2.5,
    "width": 10.333,
    "height": 3.5,
}

TITLE_FONT = "Open Sans"
SPEAKER_FONT = "Open Sans"
TITLE_MAX_PT = 36
TITLE_MIN_PT = 24
SPEAKER_MAX_PT = 21
SPEAKER_MIN_PT = 16


# -----------------------------------------------------------------------------
# DATA LOADING
# -----------------------------------------------------------------------------

def parse_is_moderator(value):
    if not value:
        return False
    return str(value).strip().lower() in ("yes", "true", "1", "y")


def load_events(csv_path):
    df = pd.read_csv(csv_path)
    events = df.groupby("event_name").agg({"brand": "first"}).reset_index()
    return events.to_dict("records")


def load_sessions_for_event(csv_path, event_name, stream=None):
    df = pd.read_csv(csv_path)
    df = df[df["event_name"] == event_name]
    if stream:
        df = df[df["stage"] == stream]
    sessions = {}
    for _, row in df.iterrows():
        session_id = str(row.get("session_id", "")).strip()
        if not session_id:
            continue
        if session_id not in sessions:
            sessions[session_id] = {
                "session_id": session_id,
                "session_title": str(row.get("session_title", "")).strip(),
                "speakers": [],
            }
        sessions[session_id]["speakers"].append({
            "name": str(row.get("speaker_name", "")).strip(),
            "job_title": str(row.get("job_title", "")).strip(),
            "company": str(row.get("company", "")).strip(),
            "is_moderator": parse_is_moderator(row.get("is_moderator", "")),
        })
    return list(sessions.values())


def load_streams_for_event(csv_path, event_name):
    df = pd.read_csv(csv_path)
    df = df[df["event_name"] == event_name]
    streams = df["stage"].dropna().unique().tolist()
    return sorted([s.strip() for s in streams if s.strip()])


def get_event_brand(csv_path, event_name):
    df = pd.read_csv(csv_path)
    event_row = df[df["event_name"] == event_name]
    if not event_row.empty:
        return str(event_row.iloc[0].get("brand", "Unknown")).strip()
    return "Unknown"


# -----------------------------------------------------------------------------
# SLIDE GENERATION
# -----------------------------------------------------------------------------

def hex_to_rgb(hex_str):
    hex_str = hex_str.lstrip("#")
    return RGBColor.from_string(hex_str)


def calculate_font_size(text, max_size, min_size, available_chars):
    if len(text) <= available_chars:
        return max_size
    ratio = available_chars / max(len(text), 1)
    return max(min_size, min(max_size, int(max_size * ratio)))


def fit_content(session, panel_height_in=CONTENT_PANEL["height"]):
    title = session["session_title"]
    speakers = session["speakers"]
    speakers = sorted(speakers, key=lambda s: (not s["is_moderator"], s["name"]))
    panel_width = CONTENT_PANEL["width"]
    title_max_chars = int(panel_width * 15 * 2)
    title_size = calculate_font_size(title, TITLE_MAX_PT, TITLE_MIN_PT, title_max_chars)
    if not speakers:
        return title_size, SPEAKER_MAX_PT
    speaker_lines = []
    for spk in speakers:
        line_len = len(spk["name"])
        if spk["job_title"]:
            line_len += len(spk["job_title"]) + 2
        if spk["company"]:
            line_len += len(spk["company"]) + 2
        if spk["is_moderator"]:
            line_len += 11
        lines_needed = max(1, (line_len + 24) // 25)
        speaker_lines.append(lines_needed)
    total_speaker_lines = sum(speaker_lines)
    speaker_area = panel_height_in * 0.6
    line_height_in = 0.25
    max_lines = int(speaker_area / line_height_in)
    if total_speaker_lines <= max_lines:
        return title_size, SPEAKER_MAX_PT
    ratio = max_lines / max(total_speaker_lines, 1)
    speaker_size = max(SPEAKER_MIN_PT, int(SPEAKER_MAX_PT * ratio))
    return title_size, speaker_size


def substitute_tokens_in_paragraph(p, session, speaker=None, scale=1.0, context=""):
    """Substitute tokens in a paragraph, preserving each run's original formatting.

    PowerPoint may split tokens across runs (e.g., "{session" in run0, "_title}" in run1).
    This function substitutes tokens into the specific run(s) they originally occupied,
    so literal separators and each token's formatting (bold name, regular job title,
    bold company) stay exactly as designed in the template.
    """
    runs = list(p.runs)
    if not runs:
        return

    run_texts = [r.text for r in runs]
    full_text = "".join(run_texts)
    char_run_index = []
    for i, t in enumerate(run_texts):
        char_run_index.extend([i] * len(t))

    print(f"DEBUG paragraph runs {context}:")
    for i, run in enumerate(runs):
        print(f"  run[{i}] = {repr(run.text)}")
    print(f"  full_text = {repr(full_text)}")

    tokens = ["{session_title}", "{speaker_name}", "{job_title}", "{company}"]
    if not any(t in full_text for t in tokens):
        print("  No tokens found, skipping")
        return

    replacements = {
        "{session_title}": session.get("session_title", ""),
        "{speaker_name}": speaker.get("name", "") if speaker else "{NAME}",
        "{job_title}": speaker.get("job_title", "") if speaker else "{JOB}",
        "{company}": speaker.get("company", "") if speaker else "{COMPANY}",
    }

    import re
    pattern = "|".join(re.escape(t) for t in tokens)
    matches = list(re.finditer(pattern, full_text))
    if not matches:
        print("  No tokens found, skipping")
        return

    new_run_text = {i: "" for i in range(len(runs))}
    cursor = 0
    for m in matches:
        start, end = m.start(), m.end()
        for pos in range(cursor, start):
            new_run_text[char_run_index[pos]] += full_text[pos]
        owner_idx = char_run_index[start]  # token inherits the formatting of its starting run
        new_run_text[owner_idx] += replacements[m.group(0)]
        cursor = end
    for pos in range(cursor, len(full_text)):
        new_run_text[char_run_index[pos]] += full_text[pos]

    print(f"  after substitution (per-run) = {[new_run_text[i] for i in range(len(runs))]}")

    for i, run in enumerate(runs):
        run.text = new_run_text[i]
        if scale < 1.0 and run.font.size and new_run_text[i]:
            run.font.size = Pt(int(run.font.size.pt * scale))


def substitute_tokens_in_shape(shape, session):
    """Substitute tokens in a shape's text frame. Returns True if this is a speaker template."""
    if not shape.has_text_frame:
        return False

    tf = shape.text_frame
    has_speaker_template = False

    for p in tf.paragraphs:
        full_text = get_paragraph_text(p)
        if "{speaker_name}" in full_text:
            has_speaker_template = True

    return has_speaker_template


def get_paragraph_text(p):
    """Get paragraph text by concatenating all runs.

    Handles PowerPoint splitting tokens across multiple runs.
    Same logic used by both generation and preview scanning.
    """
    return "".join(run.text for run in p.runs)


def strip_placeholder_identity(shape_element):
    """Remove placeholder identity from copied shapes so they render independently.

    Placeholder shapes retain references to the slide layout; stripping these
    ensures copied shapes display correctly on new slides.
    Safe no-op on shapes that aren't placeholders.
    """
    for ph in shape_element.findall('.//' + qn('p:ph')):
        ph.getparent().remove(ph)


def bake_in_inherited_geometry(shape, layout):
    """Copy placeholder's inherited position/fill/border from layout definition
    directly into shape's own XML so it renders correctly even on a layout with
    no matching placeholder to inherit from.
    """
    from copy import deepcopy
    from lxml import etree

    if not shape.is_placeholder:
        return

    idx = shape.placeholder_format.idx
    layout_ph = next(
        (lph for lph in layout.placeholders if lph.placeholder_format.idx == idx),
        None
    )
    if layout_ph is None:
        return

    sp_el = shape._element
    spPr = sp_el.find(qn('p:spPr'))
    if spPr is None:
        spPr = etree.SubElement(sp_el, qn('p:spPr'))

    layout_spPr = layout_ph._element.find(qn('p:spPr'))
    if layout_spPr is None:
        return

    # Copy transform if shape doesn't have one
    if spPr.find(qn('a:xfrm')) is None:
        layout_xfrm = layout_spPr.find(qn('a:xfrm'))
        if layout_xfrm is not None:
            spPr.insert(0, deepcopy(layout_xfrm))

    # Copy fill if shape doesn't have one
    fill_tags = ('a:solidFill', 'a:gradFill', 'a:blipFill', 'a:noFill')
    has_fill = any(spPr.find(qn(t)) is not None for t in fill_tags)
    if not has_fill:
        for fill_tag in fill_tags:
            layout_fill = layout_spPr.find(qn(fill_tag))
            if layout_fill is not None:
                spPr.append(deepcopy(layout_fill))
                break

    # Copy line/border if shape doesn't have one
    if spPr.find(qn('a:ln')) is None:
        layout_ln = layout_spPr.find(qn('a:ln'))
        if layout_ln is not None:
            spPr.append(deepcopy(layout_ln))

    # Copy geometry/prstGeom if shape doesn't have one
    if spPr.find(qn('a:prstGeom')) is None:
        layout_geom = layout_spPr.find(qn('a:prstGeom'))
        if layout_geom is not None:
            spPr.append(deepcopy(layout_geom))

    # Ensure correct OOXML child element order (xfrm, prstGeom, fill, ln)
    # PowerPoint enforces strict ordering; LibreOffice does not
    reorder_spPr_children(spPr)


def reorder_spPr_children(spPr):
    """Reorder spPr children to comply with OOXML schema sequence.

    Required order: xfrm, prstGeom, fill (any type), ln
    PowerPoint enforces this strictly; LibreOffice is lenient.
    """
    # OOXML strict order: xfrm, geometry, fill, line
    CORRECT_ORDER = [
        'a:xfrm',
        'a:custGeom',
        'a:prstGeom',
        'a:noFill',
        'a:solidFill',
        'a:gradFill',
        'a:blipFill',
        'a:pattFill',
        'a:grpFill',
        'a:ln',
    ]

    # Collect all current children by their tag
    children = {}
    for child in list(spPr):
        children[child.tag] = child
        spPr.remove(child)

    # Re-append in correct order
    for tag_name in CORRECT_ORDER:
        full_tag = qn(tag_name)
        if full_tag in children:
            spPr.append(children[full_tag])

    # Append any remaining unrecognized elements at the end
    appended_tags = {qn(t) for t in CORRECT_ORDER}
    for child_tag, child in list(children.items()):
        if child_tag not in appended_tags:
            spPr.append(child)


def scan_template_for_tokens(template_bytes):
    """Scan template PPTX to find which tokens are actually present.

    Returns a dict with token names as keys and True/False for presence.
    Tokens: session_title, speaker_name, job_title, company
    """
    import re
    from pptx import Presentation
    import tempfile
    import os

    TOKEN_PATTERNS = {
        'session_title': r'\{session_title\}',
        'speaker_name': r'\{speaker_name\}',
        'job_title': r'\{job_title\}',
        'company': r'\{company\}',
    }

    found_tokens = {
        'session_title': False,
        'speaker_name': False,
        'job_title': False,
        'company': False,
    }

    # Write to temp file and load
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        data = template_bytes.getvalue() if hasattr(template_bytes, 'getvalue') else template_bytes
        tmp.write(data)
        tmp_path = tmp.name

    try:
        prs = Presentation(tmp_path)
        # Check all slides for tokens
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for p in shape.text_frame.paragraphs:
                    para_text = get_paragraph_text(p)
                    for token_name, pattern in TOKEN_PATTERNS.items():
                        if re.search(pattern, para_text):
                            found_tokens[token_name] = True
    finally:
        os.unlink(tmp_path)

    return found_tokens


def extract_slide_confirmation(pptx_bytes, template_bytes=None):
    """Extract text confirmation from generated PPTX showing substituted values with per-token status.

    Only checks and displays tokens that exist in the original template.

    Args:
        pptx_bytes: The generated output PPTX
        template_bytes: The original template PPTX (optional, for token scanning)

    Returns a list of dicts with 'slide_num', 'session_title', 'session_title_resolved',
    'session_title_in_template', 'speakers'.
    Each speaker has per-token resolution status and template presence flags.
    """
    import re
    from pptx import Presentation
    import tempfile
    import os

    TOKEN_PATTERNS = {
        'session_title': r'\{session_title\}',
        'speaker_name': r'\{speaker_name\}',
        'job_title': r'\{job_title\}',
        'company': r'\{company\}',
    }

    # First, scan template to find which tokens exist (if template_bytes provided)
    tokens_in_template = {
        'session_title': template_bytes is None,  # If no template provided, assume all exist
        'speaker_name': template_bytes is None,
        'job_title': template_bytes is None,
        'company': template_bytes is None,
    }
    if template_bytes is not None:
        tokens_in_template = scan_template_for_tokens(template_bytes)
        print(f"DEBUG scan result: {tokens_in_template}")

    def check_token_in_text(text, token_pattern):
        """Check if token pattern still exists in text (unresolved)."""
        return bool(re.search(token_pattern, text))

    def is_substantial_text(text):
        """Check if text is substantial (not placeholder like '.')."""
        return len(text) > 2 and any(c.isalnum() for c in text)

    slides = []

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        data = pptx_bytes.getvalue() if hasattr(pptx_bytes, 'getvalue') else pptx_bytes
        tmp.write(data)
        tmp_path = tmp.name

    try:
        prs = Presentation(tmp_path)
        for slide_idx, slide in enumerate(prs.slides, 1):
            session_title_value = None
            session_title_resolved = False
            session_title_in_template = tokens_in_template.get('session_title', False)
            speakers = []
            current_speaker = None

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                tf = shape.text_frame

                for p in tf.paragraphs:
                    para_text = get_paragraph_text(p).strip()

                    if not is_substantial_text(para_text):
                        continue

                    # Check for unresolved tokens (placeholder still present)
                    has_session_title = check_token_in_text(para_text, TOKEN_PATTERNS['session_title'])
                    has_speaker_name = check_token_in_text(para_text, TOKEN_PATTERNS['speaker_name'])
                    has_job_title = check_token_in_text(para_text, TOKEN_PATTERNS['job_title'])
                    has_company = check_token_in_text(para_text, TOKEN_PATTERNS['company'])

                    # Classify paragraph based on what tokens are present
                    if has_speaker_name:
                        # This is a speaker line
                        parts = para_text.replace('|', ',').split(',')
                        name_part = parts[0].strip() if len(parts) > 0 else ""
                        job_part = parts[1].strip() if len(parts) > 1 else ""
                        company_part = parts[2].strip() if len(parts) > 2 else ""

                        speakers.append({
                            "name": name_part,
                            "name_resolved": not has_speaker_name,
                            "name_in_template": tokens_in_template.get('speaker_name', False),
                            "job_title": job_part,
                            "job_title_resolved": not has_job_title,
                            "job_title_in_template": tokens_in_template.get('job_title', False),
                            "company": company_part,
                            "company_resolved": not has_company,
                            "company_in_template": tokens_in_template.get('company', False),
                        })
                    elif has_session_title and session_title_value is None:
                        # This is the session title
                        session_title_value = para_text
                        session_title_resolved = not has_session_title

            slides.append({
                "slide_num": slide_idx,
                "session_title": session_title_value or "(no title found)",
                "session_title_resolved": session_title_resolved,
                "session_title_in_template": session_title_in_template,
                "speakers": speakers
            })
    finally:
        os.unlink(tmp_path)

    return slides



def generate_slides_from_template(template_bytes, sessions):
    """Generate one slide per session using a PPTX template with token replacement.

    Two-phase process, strictly ordered:
      Phase 1: capture pristine copies of the template's shapes, then create N
               slides (N = len(sessions)) with NO substitution applied to any of them.
      Phase 2: substitute each slide's own session data independently, only after
               all N slides already exist as separate objects.

    Tokens: {session_title}, {speaker_name}, {job_title}, {company}
    """
    output = io.BytesIO()

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp_template:
        tmp_template.write(template_bytes)
        tmp_template_path = tmp_template.name

    try:
        template_prs = Presentation(tmp_template_path)
        original_slide = template_prs.slides[0]
        original_layout = original_slide.slide_layout

        # --- Bake inherited geometry BEFORE capturing shapes ---
        # Panel shape inherits position/fill from layout; bake it in so it works
        # on Blank layouts that don't have the same placeholder definitions.
        for shape in original_slide.shapes:
            bake_in_inherited_geometry(shape, original_layout)
        print(f"DEBUG: Baked inherited geometry from layout {original_layout.name!r}")

        # --- Capture pristine shape elements BEFORE anything is touched ---
        pristine_shape_elements = [deepcopy(shape._element) for shape in original_slide.shapes]
        print(f"DEBUG: Captured {len(pristine_shape_elements)} pristine shape(s) from original slide")

        # --- Find a blank layout (no title/body placeholders) for new slides ---
        blank_layout = None
        for layout in template_prs.slide_masters[0].slide_layouts:
            if layout.name.strip().lower() == "blank":
                blank_layout = layout
                break
        if blank_layout is None:
            blank_layout = min(
                template_prs.slide_masters[0].slide_layouts,
                key=lambda l: len(l.placeholders)
            )
        print(f"DEBUG: Using layout {blank_layout.name!r} for new slides")

        # ============ PHASE 1: create all N slides, zero substitution ============
        n = len(sessions)
        all_slides = [original_slide]  # slide 1 already exists and is still pristine

        for i in range(1, n):
            new_slide = template_prs.slides.add_slide(blank_layout)
            for shp in list(new_slide.shapes):
                shp._element.getparent().remove(shp._element)
            for el in pristine_shape_elements:
                el_copy = deepcopy(el)
                strip_placeholder_identity(el_copy)
                new_slide.shapes._spTree.append(el_copy)
            for rel_id, rel in original_slide.part.rels.items():
                if "notesSlide" not in rel.reltype and "slideLayout" not in rel.reltype:
                    if "image" in rel.reltype:
                        new_slide.part.rels.add_relationship(rel.reltype, rel.target_part, rel_id)
            all_slides.append(new_slide)
            print(f"DEBUG: Created slide {i + 1} of {n} (pristine, no substitution yet)")

        print(f"DEBUG: Phase 1 complete — created {len(all_slides)} slides, zero substitution done")

        # ============ PHASE 2: substitute each slide's own session data ============
        for i, (slide, session) in enumerate(zip(all_slides, sessions), start=1):
            print(f"DEBUG: Phase 2 — substituting session {i} of {n}: {session['session_title']}")
            speakers = session["speakers"]
            speaker_count = len(speakers)

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                tf = shape.text_frame

                speaker_template_paras = []
                title_paras = []
                for p in tf.paragraphs:
                    full_text = get_paragraph_text(p)
                    if "{speaker_name}" in full_text:
                        speaker_template_paras.append(p)
                    elif "{session_title}" in full_text:
                        title_paras.append(p)

                if speaker_template_paras and speaker_count > 0:
                    template_para = speaker_template_paras[0]

                    # Pristine copy captured BEFORE any substitution touches it —
                    # this is what fixes the "same speaker repeated" bug.
                    pristine_speaker_el = deepcopy(template_para._element)

                    total_lines = speaker_count
                    available_height = shape.height
                    estimated_line_height = 360000
                    available_lines = max(3, available_height / estimated_line_height)
                    scale = max(0.6, available_lines / total_lines) if total_lines > available_lines else 1.0

                    substitute_tokens_in_paragraph(
                        template_para, session, speakers[0], scale,
                        context=f"S{i}-SPEAKER-1/{speaker_count}"
                    )

                    for idx, speaker in enumerate(speakers[1:], start=2):
                        new_p_el = deepcopy(pristine_speaker_el)
                        tf._element.append(new_p_el)
                        new_p = tf.paragraphs[-1]
                        substitute_tokens_in_paragraph(
                            new_p, session, speaker, scale,
                            context=f"S{i}-SPEAKER-{idx}/{speaker_count}"
                        )

                for p in title_paras:
                    substitute_tokens_in_paragraph(
                        p, session, None, 1.0, context=f"S{i}-TITLE"
                    )

        template_prs.save(output)
        output.seek(0)

    finally:
        os.unlink(tmp_template_path)

    return output, len(sessions)



def generate_slides(csv_path, event_name, stream, template_image_bytes, text_color):
    sessions = load_sessions_for_event(csv_path, event_name, stream)
    prs = Presentation()
    prs.slide_width = Emu(int(Inches(SLIDE_WIDTH_IN)))
    prs.slide_height = Emu(int(Inches(SLIDE_HEIGHT_IN)))
    for session in sessions:
        add_slide_with_template(
            prs, session,
            template_image_bytes=template_image_bytes,
            text_color=text_color,
        )
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output, len(sessions)


# -----------------------------------------------------------------------------
# PIL PREVIEW RENDERING
# -----------------------------------------------------------------------------

def hex_to_rgb_tuple(hex_str):
    hex_str = hex_str.lstrip("#")
    return tuple(int(hex_str[i:i+2], 16) for i in (0, 2, 4))


def render_slide_preview_pil(target_width_px, session, template_image_bytes, text_color):
    """Render a slide preview using PIL Image with pre-designed template.

    Just overlays dynamic text (session title, speaker list) onto the template image.
    """
    target_height_px = int(target_width_px / SLIDE_ASPECT)

    # Template as background
    if template_image_bytes:
        try:
            bg = Image.open(io.BytesIO(template_image_bytes))
            bg = bg.resize((target_width_px, target_height_px), Image.LANCZOS)
        except:
            bg_color = hex_to_rgb_tuple(COLORS["default_bg"])
            bg = Image.new("RGB", (target_width_px, target_height_px), bg_color)
    else:
        bg_color = hex_to_rgb_tuple(COLORS["default_bg"])
        bg = Image.new("RGB", (target_width_px, target_height_px), bg_color)

    if bg.mode != 'RGBA':
        bg = bg.convert('RGBA')

    draw = ImageDraw.Draw(bg, 'RGBA')

    # Load bundled fonts
    script_dir = Path(__file__).parent
    fonts_dir = script_dir / "fonts"
    try:
        title_font = ImageFont.truetype(str(fonts_dir / "OpenSans-Bold.ttf"), int(target_width_px * 0.04))
        speaker_font_regular = ImageFont.truetype(str(fonts_dir / "OpenSans-Regular.ttf"), int(target_width_px * 0.023))
        speaker_font_bold = ImageFont.truetype(str(fonts_dir / "OpenSans-Bold.ttf"), int(target_width_px * 0.023))
    except Exception as e:
        st.warning(f"Bundled fonts not found at {fonts_dir}, using default font: {e}")
        title_font = ImageFont.load_default()
        speaker_font_regular = ImageFont.load_default()
        speaker_font_bold = ImageFont.load_default()

    # Content panel positions
    panel_left = int((CONTENT_PANEL["left"] / SLIDE_WIDTH_IN) * target_width_px)
    panel_top = int((CONTENT_PANEL["top"] / SLIDE_HEIGHT_IN) * target_height_px)
    panel_right = int(((CONTENT_PANEL["left"] + CONTENT_PANEL["width"]) / SLIDE_WIDTH_IN) * target_width_px)

    # Session title (bold)
    title_x = panel_left + int(target_width_px * 0.025)
    title_y = panel_top + int(target_height_px * 0.08)
    title = session.get("session_title", "")
    text_color_tuple = hex_to_rgb_tuple(text_color or COLORS["title_text_dark"])

    # Simple text wrapping for title
    max_title_width = panel_right - panel_left - int(target_width_px * 0.05)
    words = title.split()
    lines = []
    current_line = ""
    for word in words:
        test_line = current_line + " " + word if current_line else word
        bbox = draw.textbbox((0, 0), test_line, font=title_font)
        if bbox[2] - bbox[0] <= max_title_width:
            current_line = test_line
        else:
            if current_line:
                lines.append(current_line)
            current_line = word
    if current_line:
        lines.append(current_line)

    # Draw wrapped title
    line_height = int(target_height_px * 0.06)
    for i, line in enumerate(lines[:2]):  # Max 2 lines
        draw.text((title_x, title_y + i * line_height), line, font=title_font, fill=text_color_tuple)

    # Speakers with proper bolding
    speakers = sorted(session.get("speakers", []), key=lambda s: (not s["is_moderator"], s["name"]))
    speaker_y = panel_top + int(target_height_px * 0.28)
    line_spacing = int(target_height_px * 0.055)
    space_width = int(target_width_px * 0.008)

    for speaker in speakers:
        current_x = title_x

        # Moderator label (not bold)
        if speaker["is_moderator"]:
            mod_text = "Moderator: "
            draw.text((current_x, speaker_y), mod_text, font=speaker_font_regular, fill=text_color_tuple)
            bbox = draw.textbbox((current_x, speaker_y), mod_text, font=speaker_font_regular)
            current_x += (bbox[2] - bbox[0]) + space_width

        # Name (bold)
        name = speaker["name"]
        draw.text((current_x, speaker_y), name, font=speaker_font_bold, fill=text_color_tuple)
        bbox = draw.textbbox((current_x, speaker_y), name, font=speaker_font_bold)
        current_x += (bbox[2] - bbox[0])

        # Job title (not bold)
        if speaker["job_title"]:
            comma = ", "
            draw.text((current_x, speaker_y), comma, font=speaker_font_regular, fill=text_color_tuple)
            bbox = draw.textbbox((current_x, speaker_y), comma, font=speaker_font_regular)
            current_x += (bbox[2] - bbox[0])
            draw.text((current_x, speaker_y), speaker["job_title"], font=speaker_font_regular, fill=text_color_tuple)
            bbox = draw.textbbox((current_x, speaker_y), speaker["job_title"], font=speaker_font_regular)
            current_x += (bbox[2] - bbox[0])

        # Company (bold)
        if speaker["company"]:
            comma = ", "
            draw.text((current_x, speaker_y), comma, font=speaker_font_regular, fill=text_color_tuple)
            bbox = draw.textbbox((current_x, speaker_y), comma, font=speaker_font_regular)
            current_x += (bbox[2] - bbox[0])
            draw.text((current_x, speaker_y), speaker["company"], font=speaker_font_bold, fill=text_color_tuple)

        speaker_y += line_spacing

    return bg


# -----------------------------------------------------------------------------
# STREAMLIT INTERFACE
# -----------------------------------------------------------------------------

def main():
    st.set_page_config(
        page_title="Holding Slide Generator | Informa",
        page_icon="Informa_Orbit_RGB.png",
        layout="wide",
        initial_sidebar_state="collapsed",
    )

    # Informa logo and title header (stacked vertically)
    informa_logo_path = Path(__file__).parent / "Informa Logo" / "Informa_Logo_OneLine_Graduated_Indigo_RGB.png"
    if informa_logo_path.exists():
        st.image(str(informa_logo_path), width=120)

    # Title with brand color and subtitle (Fibonacci spacing: 13px, 21px)
    st.markdown(
        f"""
        <h1 style='font-size: 2.5rem; margin-bottom: 13px; margin-top: 8px; color: #{COLORS["cta_bg"]};'>
            Holding slide generator
        </h1>
        <p style='color: #6B7280; font-size: 1rem; margin-bottom: 21px;'>
            Generate professional holding slides for your event sessions
        </p>
        """,
        unsafe_allow_html=True
    )

    # Initialize session state
    if "csv_path" not in st.session_state:
        st.session_state.csv_path = "sample_sessions.csv"

    # Initialize control values in session state
    if "selected_event" not in st.session_state:
        st.session_state.selected_event = None
    if "selected_stream" not in st.session_state:
        st.session_state.selected_stream = None

# Controls section — split into two columns: controls (left) and example (right)
    csv_path = "sample_sessions.csv"
    selected_event = None
    selected_stream = None
    template_bytes = None
    text_color = "#" + COLORS["title_text_dark"]

    # Create two columns — narrower left for controls, wider right for example
    main_cols = st.columns([26, 74])

    # LEFT COLUMN: Controls
    with main_cols[0]:
        # Collapsed expander for template format instructions (Streamlit constrains width to column)
        example_image_path = Path(__file__).parent / "Example Markup.png"
        with st.expander("How to format your template", expanded=False):
            st.markdown("""
            - Tokens: `{session_title}`, `{speaker_name}`, `{job_title}`, `{company}`}
            - Font, weight, color carry over from how you style the tokens
            - One speaker line handles panels too — repeats automatically per speaker
            """)
            # Example template image inside expander (fits expander width)
            if example_image_path.exists():
                st.image(str(example_image_path), use_container_width=True, output_format="PNG")
            else:
                st.caption("Example image not found")

        # Configuration card — Event, Stream, Template grouped together (Proximity principle)
        with st.container(border=True):
            # Event selection (13px gap after)
            st.markdown("**Event**")
            try:
                events = load_events(csv_path)
                if events:
                    event_names = [e["event_name"] for e in events]
                    selected_event = st.selectbox(
                        "",
                        event_names,
                        key="event_select",
                        label_visibility="collapsed"
                    )
                    st.session_state.selected_event = selected_event
                else:
                    st.warning("No events found")
                    selected_event = None
            except Exception as e:
                st.error(f"Could not load events: {e}")
                selected_event = None

            # Stream selection (21px gap — section break within card)
            st.markdown("<div style='margin-top: 21px;'></div>", unsafe_allow_html=True)
            st.markdown("**Stream**")
            if selected_event:
                try:
                    streams = load_streams_for_event(csv_path, selected_event)
                    if streams:
                        selected_stream = st.selectbox(
                            "",
                            streams,
                            key="stream_select",
                            label_visibility="collapsed"
                        )
                        st.session_state.selected_stream = selected_stream
                    else:
                        st.warning("No streams found")
                        selected_stream = None
                except Exception as e:
                    st.error(f"Could not load streams: {e}")
                    selected_stream = None
            else:
                st.selectbox("", ["Select event first"], disabled=True, label_visibility="collapsed")
                selected_stream = None

            # Template upload (21px gap)
            st.markdown("<div style='margin-top: 21px;'></div>", unsafe_allow_html=True)
            st.markdown("**Template**")
            template_file = st.file_uploader(
                "Upload PPTX",
                type=["pptx"],
                help="Upload a PPTX template with {session_title}, {speaker_name}, {job_title}, {company} tokens",
                label_visibility="collapsed"
            )
            template_bytes = template_file.getvalue() if template_file else None

        # Generate button — separate from config card (34px gap, visual distinction)
        st.markdown("<div style='margin-top: 34px;'></div>", unsafe_allow_html=True)
        if selected_event and selected_stream:
            st.markdown("""
            <style>
                .stButton>button {
                    background-color: #002244 !important;
                    border-color: #002244 !important;
                }
                .stButton>button:hover {
                    background-color: #003CB2 !important;
                    border-color: #003CB2 !important;
                }
                .stButton>button:focus {
                    background-color: #003CB2 !important;
                    border-color: #003CB2 !important;
                }
                .stButton>button:active {
                    background-color: #003CB2 !important;
                    border-color: #003CB2 !important;
                }
            </style>
            """, unsafe_allow_html=True)
            generate_clicked = st.button(
                "Generate slides",
                type="primary",
                use_container_width=True,
            )

            if generate_clicked:
                with st.spinner("Generating slides..."):
                    try:
                        sessions = load_sessions_for_event(csv_path, selected_event, selected_stream)
                        output_bytes, num_slides = generate_slides_from_template(
                            template_bytes, sessions
                        )

                        # Auto-trigger download via hidden link
                        filename = f"{selected_event.replace(' ', '_')}_{selected_stream.replace(' ', '_')}_holding_slides.pptx"
                        b64 = base64.b64encode(output_bytes.getvalue()).decode()
                        st_html(f"""
                            <a id="auto-dl" href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64}" download="{filename}"></a>
                            <script>document.getElementById('auto-dl').click();</script>
                        """, height=0)

                        # Generation complete - download triggered above
                    except Exception as e:
                        st.error(f"Error: {e}")

    # RIGHT COLUMN: Empty (all content moved to left column expander)
    with main_cols[1]:
        pass

if __name__ == "__main__":
    main()
